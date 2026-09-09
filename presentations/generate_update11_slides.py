"""Generate FYP Update 11 Presentation — Auditing My Own Results

Update 10 closed by naming cross-sectional as the last open thread.
Going back to close it out, the results turned out not to reproduce,
and investigating that surfaced six defects in the testing methodology
itself. This deck is about that audit and the corrected re-run.

Card text pulls the corrected numbers from results/ at build time so
the deck cannot drift from the study output.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import glob
import sys

import pandas as pd

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0xED, 0xE8, 0xE0)
BG_CARD = RGBColor(0xE0, 0xDB, 0xD3)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
BRACKET = RGBColor(0x33, 0x33, 0x2D)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
DIM = RGBColor(0x6B, 0x6B, 0x63)
MID = RGBColor(0x55, 0x55, 0x50)

CHART_DIR = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent
RESULTS = Path(__file__).resolve().parent.parent / 'results'


def latest(pattern):
    files = sorted(glob.glob(str(RESULTS / pattern)))
    return files[-1] if files else None


corrected_path = latest('cross_sectional_corrected_*.csv')
if corrected_path:
    cdf = pd.read_csv(corrected_path)
    N_ROBUST = int(cdf['robust'].sum())
    N_TOTAL = len(cdf)
    CORRECTED_LINE = f"{N_ROBUST} of {N_TOTAL} market-strategy combinations validated."
else:
    N_ROBUST, N_TOTAL = 0, 4
    CORRECTED_LINE = "(corrected re-run pending)"


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False,
             alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT,
             font_name="Calibri", space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bracket_tl(slide, left, top, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_bracket_br(slide, right, bottom, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - thickness), Inches(bottom - size), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - size), Inches(bottom - thickness), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_divider(slide, left, top, width):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    s.fill.solid(); s.fill.fore_color.rgb = BRACKET; s.line.fill.background()


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    s.fill.solid(); s.fill.fore_color.rgb = fill_color
    s.line.color.rgb = RGBColor(0xCC, 0xC7, 0xBF); s.line.width = Pt(1)
    return s


# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 1.0, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.5, size=2.0, thickness=0.15)
add_text(slide, 2.5, 2.6, 8.5, 1.2, "FYP UPDATE 11", font_size=44, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 3.8, 8.5, 0.8, "Auditing My Own Results", font_size=20, color=MID, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 5.0, 8.5, 0.5, "Jackson Seow  •  FYP 2025/2026", font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 2: THE RESULT THAT WOULDN'T REPRODUCE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "The Result That Wouldn't Reproduce", font_size=28, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

img = CHART_DIR / 'chart_u11_repro.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(2.1), Inches(1.5), width=Inches(9.1))

add_card(slide, 0.8, 6.15, 11.8, 1.15)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5,
              "Testing cross-sectional reversal on the S&P 100 looked promising at first. Re-running it to confirm gave a different answer.",
              font_size=14, color=DARK, bold=True)
add_para(tf, "Same strategy, same stocks, same settings. Only the day it was run changed, which suggested the problem was in the testing code rather than the strategy.",
         font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 3: WHAT THE AUDIT FOUND ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "Six Defects in My Own Testing", font_size=28, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 3.8, 4.9)
tf = add_text(slide, 1.05, 1.7, 3.3, 0.5, "Inflated the results", font_size=15, color=ACCENT_RED, bold=True)
add_para(tf, "")
add_para(tf, "1.  Training data leaked into", font_size=12.5, color=DARK)
add_para(tf, "     testing. Each window's", font_size=12.5, color=DARK)
add_para(tf, "     boundary day sat in both", font_size=12.5, color=DARK)
add_para(tf, "     sets: 2% of the test data", font_size=12.5, color=DARK)
add_para(tf, "     under one setting, 4.9%", font_size=12.5, color=DARK)
add_para(tf, "     under the other.", font_size=12.5, color=DARK)
add_para(tf, "")
add_para(tf, "The strategy was partly being", font_size=12.5, color=MID)
add_para(tf, "marked on work it had already", font_size=12.5, color=MID)
add_para(tf, "seen.", font_size=12.5, color=MID)

add_card(slide, 4.8, 1.5, 3.8, 4.9)
tf = add_text(slide, 5.05, 1.7, 3.3, 0.5, "Depressed the results", font_size=15, color=ACCENT_GREEN, bold=True)
add_para(tf, "")
add_para(tf, "2.  Hong Kong broker fees", font_size=12.5, color=DARK)
add_para(tf, "     charged on every US trade,", font_size=12.5, color=DARK)
add_para(tf, "     about 0.32% per round trip", font_size=12.5, color=DARK)
add_para(tf, "     that US trading would not", font_size=12.5, color=DARK)
add_para(tf, "     actually incur.", font_size=12.5, color=DARK)
add_para(tf, "")
add_para(tf, "3.  Test periods running past", font_size=12.5, color=DARK)
add_para(tf, "     the end of the data scored", font_size=12.5, color=DARK)
add_para(tf, "     a flat 0%, and that counted", font_size=12.5, color=DARK)
add_para(tf, "     as a losing period.", font_size=12.5, color=DARK)

add_card(slide, 8.8, 1.5, 3.8, 4.9, fill_color=RGBColor(0xE4, 0xDE, 0xD2))
tf = add_text(slide, 9.05, 1.7, 3.3, 0.5, "Made them unreliable", font_size=15, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "4.  The answer changed run to", font_size=12.5, color=DARK)
add_para(tf, "     run, from an ordering that", font_size=12.5, color=DARK)
add_para(tf, "     Python randomises.", font_size=12.5, color=DARK)
add_para(tf, "")
add_para(tf, "5.  Windows anchored to today's", font_size=12.5, color=DARK)
add_para(tf, "     date, so the test period slid", font_size=12.5, color=DARK)
add_para(tf, "     forward daily.", font_size=12.5, color=DARK)
add_para(tf, "")
add_para(tf, "6.  Fixed 100-share lots, so the", font_size=12.5, color=DARK)
add_para(tf, "     money committed tracked", font_size=12.5, color=DARK)
add_para(tf, "     share price, not any decision.", font_size=12.5, color=DARK)

add_card(slide, 0.8, 6.55, 11.8, 0.7)
add_text(slide, 1.1, 6.65, 11.3, 0.5,
         "All six are fixed. The full test suite passes, and the same study now returns an identical number on every run.",
         font_size=13, color=MID)


# ==================== SLIDE 4: THE MOST SERIOUS ONE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "The Most Serious One", font_size=28, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

img = CHART_DIR / 'chart_u11_codefix.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(1.65), Inches(1.45), width=Inches(10.0))

img = CHART_DIR / 'chart_u11_determinism.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(3.32), Inches(3.15), width=Inches(6.7))

add_card(slide, 0.8, 6.15, 11.8, 1.15)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5,
              "The basket was held in a Python set, whose ordering changes between runs, and an unaffordable stock is skipped.",
              font_size=14, color=DARK, bold=True)
add_para(tf, "So the answer depended on ordering. The two levels are not comparable, since fixing the sizing changed the portfolio entirely; "
             "what matters is that the three bars on the right are identical.",
         font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 5: CORRECTED RESULTS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "Re-Run, Corrected", font_size=28, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

img = CHART_DIR / 'chart_u11_corrected.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(1.77), Inches(1.5), width=Inches(9.8))

add_card(slide, 0.8, 6.15, 11.8, 1.15)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5,
              f"Both directions, both market indices, all six defects corrected.  {CORRECTED_LINE}",
              font_size=14, color=DARK, bold=True)
add_para(tf, "US results are run with zero trading costs, as a deliberate best case: there are no live US fills to calibrate a fee from, "
             "and a null result that survives zero costs is stronger than one assuming an invented rate.",
         font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 6: BUT ONLY WITH NO FEES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "It Only Survives With No Trading Fees", font_size=28, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

img = CHART_DIR / 'chart_u11_costsens.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(1.87), Inches(1.5), width=Inches(9.6))

add_card(slide, 0.8, 6.15, 11.8, 1.15)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5,
              "Charging even 0.03% per side, about the cheapest US retail execution available, is enough to break it.",
              font_size=14, color=DARK, bold=True)
add_para(tf, "The passing version also leans on two windows out of nine: without those two, the remaining seven average -1.64%. "
             "It clears the bar, but not by enough to rely on.",
         font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 7: SUMMARY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 0.7, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.8, size=2.0, thickness=0.15)

add_text(slide, 2.5, 1.0, 8.5, 0.7, "SUMMARY", font_size=32, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

tf = add_text(slide, 2.0, 2.05, 9.5, 0.5, "✅   Went back to verify a result, and found it did not reproduce", font_size=16.5, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Audited my own testing code: six defects, one inflating, two depressing, three destabilising", font_size=16.5, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Fixed all six; the same study now returns an identical number on every run", font_size=16.5, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Re-ran cleanly: one combination finally validated, the first in this project", font_size=16.5, color=DARK)
add_para(tf, "")
add_para(tf, "✅   But it holds only at zero fees, and rests on 2 windows out of 9", font_size=16.5, color=DARK)

add_card(slide, 2.0, 5.25, 9.5, 1.35)
tf = add_text(slide, 2.3, 5.3, 9, 0.5, "🎯  Next: Calibrating Real US Fees, Then Writing It Up", font_size=18, color=ACCENT_GREEN, bold=True)
add_para(tf, "The verdict turns on the fee assumption, so the next step is measuring it from live fills instead of estimating.", font_size=14, color=MID)
add_para(tf, "One combination passed, and the same checks gave reason to doubt it.",
         font_size=14, color=MID, space_before=Pt(2))

add_text(slide, 2.5, 6.7, 8.5, 0.6, "Thank You", font_size=26, color=BRACKET, bold=True, alignment=PP_ALIGN.CENTER)


output_path = "/Users/jacksonetherchainstake/FYP/Documents/FYP_Update_11.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
print(f"  corrected results: {CORRECTED_LINE}")

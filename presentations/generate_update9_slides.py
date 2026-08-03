"""Generate FYP Update 9 Presentation — The Verdict: Does a Smarter Model
Beat the Simple Rule?

Matches the FYP theme (see generate_update5_slides.py). Covers the full
strategy research arc: rule vs HMM vs ML classifier vs cross-sectional,
the two combination attempts that didn't rescue it, and the robust
validation study extended across 3x the data and 8 new stocks, landing
on a 0/57 null result.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import sys

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0xED, 0xE8, 0xE0)
BG_CARD = RGBColor(0xE0, 0xDB, 0xD3)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
BRACKET = RGBColor(0x33, 0x33, 0x2D)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_AMBER = RGBColor(0xE6, 0x8A, 0x00)
DIM = RGBColor(0x6B, 0x6B, 0x63)
MID = RGBColor(0x55, 0x55, 0x50)

CHART_DIR = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bracket_tl(slide, left, top, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_bracket_br(slide, right, bottom, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - thickness), Inches(bottom - size), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - size), Inches(bottom - thickness), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = BRACKET; shape.line.fill.background()


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xCC, 0xC7, 0xBF); shape.line.width = Pt(1)
    return shape


# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 1.0, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.5, size=2.0, thickness=0.15)
add_text(slide, 2.5, 2.6, 8.5, 1.2, "FYP UPDATE 9", font_size=44, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 3.8, 8.5, 0.8, "The Verdict: Does a Smarter Model Beat the Simple Rule?", font_size=20, color=MID, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 5.0, 8.5, 0.5, "Jackson Seow  •  FYP 2025/2026", font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 2: FOUR APPROACHES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "Four Approaches, Tested Head to Head", font_size=28, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.4, 4.5)
tf = add_text(slide, 1.1, 1.7, 4.8, 0.5, "What Was Tested", font_size=19, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "•  Rule-based regime switch (a hand-tuned", font_size=14, color=DARK)
add_para(tf, "   efficiency-ratio threshold)", font_size=14, color=DARK)
add_para(tf, "•  HMM regime switch (a model that learns", font_size=14, color=DARK)
add_para(tf, "   market conditions on its own)", font_size=14, color=DARK)
add_para(tf, "•  A logistic regression classifier predicting", font_size=14, color=DARK)
add_para(tf, "   next-candle direction", font_size=14, color=DARK)
add_para(tf, "•  Cross-sectional reversal, ranking stocks", font_size=14, color=DARK)
add_para(tf, "   against each other instead of against time", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "Then two attempts to combine the best of both", font_size=14, color=DARK)
add_para(tf, "worlds: an adaptive selector choosing rule or", font_size=14, color=DARK)
add_para(tf, "HMM per window, and a fixed 50/50 blend.", font_size=14, color=DARK)

img = CHART_DIR / 'chart_u9_combinations.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(6.5), Inches(1.7), width=Inches(6.2))
add_text(slide, 6.5, 5.55, 6.2, 0.5, "Mean out-of-sample return, same 11 stocks, 9 windows each",
         font_size=13, color=DIM)

add_card(slide, 0.8, 6.15, 11.8, 1.0)
tf = add_text(slide, 1.1, 6.3, 11.3, 0.5, "Neither combination attempt fixed the underlying problem.", font_size=14, color=DARK, bold=True)
add_para(tf, "The adaptive selector picked based on training performance, but training performance didn't predict the out-of-sample winner.", font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 3: HOW SURE ARE WE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "How Sure Are We There's No Edge?", font_size=28, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

img = CHART_DIR / 'chart_u9_funnel.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(1.3), Inches(1.5), width=Inches(10.7))

add_card(slide, 0.8, 6.15, 11.8, 1.15)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5, "A single test on 11 stocks looked promising for two pairs. Requiring two independent tests to", font_size=14, color=DARK, bold=True)
add_para(tf, "agree dropped that to zero, and it stayed at zero after tripling the data and testing 8 new stocks in different sectors. Combined: 0 out of 57.", font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 4: SUMMARY & NEXT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 0.7, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.8, size=2.0, thickness=0.15)

add_text(slide, 2.5, 1.1, 8.5, 0.7, "SUMMARY", font_size=32, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

tf = add_text(slide, 2.0, 2.1, 9.5, 0.5, "✅   Four strategy families and two combination methods tested head to head", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Validation extended to 3 years of data and 8 stocks in new sectors, same result", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Corrected an earlier claim of mine after testing it directly: the model's most flexible", font_size=17, color=DARK)
add_para(tf, "       setting doesn't help either", font_size=17, color=DARK, space_before=Pt(0))
add_para(tf, "")
add_para(tf, "✅   Take-profit now live on the roster, and every study is browsable from the dashboard", font_size=17, color=DARK)

add_card(slide, 2.0, 5.15, 9.5, 1.3)
tf = add_text(slide, 2.3, 5.3, 9, 0.5, "🎯  Next: Writing This Up Honestly", font_size=18, color=ACCENT_GREEN, bold=True)
add_para(tf, "A validated null result is still a result. The engineering built to reach it, live deployment,", font_size=14, color=MID)
add_para(tf, "risk management, and rigorous testing, is what the final report will stand on.", font_size=14, color=MID, space_before=Pt(2))

add_text(slide, 2.5, 6.7, 8.5, 0.6, "Thank You", font_size=26, color=BRACKET, bold=True, alignment=PP_ALIGN.CENTER)


output_path = "/Users/jacksonetherchainstake/FYP/Documents/FYP_Update_9.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")

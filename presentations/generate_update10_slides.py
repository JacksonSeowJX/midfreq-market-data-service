"""Generate FYP Update 10 Presentation — Can Risk Management Alone Beat
Buy-and-Hold?

Matches the FYP theme (see generate_update9_slides.py). Direct
follow-up to Update 9's closing line, which named this exact question
as the one thing still open before writing up the report: whether the
stop-loss / trailing-stop infrastructure already built into the
framework can improve on unconditional buy-and-hold, without trying to
time entries or exits at all.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import sys

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0xED, 0xE8, 0xE0)
BG_CARD = RGBColor(0xE0, 0xDB, 0xD3)
DARK = RGBColor(0x2D, 0x2D, 0x2D)
BRACKET = RGBColor(0x33, 0x33, 0x2D)
ACCENT_GREEN = RGBColor(0x2E, 0x7D, 0x32)
ACCENT_RED = RGBColor(0xC6, 0x28, 0x28)
ACCENT_AMBER = RGBColor(0xE6, 0x8A, 0x00)
DIM = RGBColor(0x6B, 0x6B, 0x63)
MID = RGBColor(0x55, 0x55, 0x50)

CHART_DIR = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent


def set_slide_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_para(tf, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", space_before=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    return p


def add_bracket_tl(slide, left, top, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_bracket_br(slide, right, bottom, size=1.2, thickness=0.12):
    v = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - thickness), Inches(bottom - size), Inches(thickness), Inches(size))
    v.fill.solid(); v.fill.fore_color.rgb = BRACKET; v.line.fill.background()
    h = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(right - size), Inches(bottom - thickness), Inches(size), Inches(thickness))
    h.fill.solid(); h.fill.fore_color.rgb = BRACKET; h.line.fill.background()


def add_divider(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(2))
    shape.fill.solid(); shape.fill.fore_color.rgb = BRACKET; shape.line.fill.background()


def add_card(slide, left, top, width, height, fill_color=BG_CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = RGBColor(0xCC, 0xC7, 0xBF); shape.line.width = Pt(1)
    return shape


# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 1.0, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.5, size=2.0, thickness=0.15)
add_text(slide, 2.5, 2.6, 8.5, 1.2, "FYP UPDATE 10", font_size=44, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 3.8, 8.5, 0.8, "Can Risk Management Alone Beat Buy-and-Hold?", font_size=20, color=MID, alignment=PP_ALIGN.CENTER)
add_text(slide, 2.5, 5.0, 8.5, 0.5, "Jackson Seow  •  FYP 2025/2026", font_size=16, color=DIM, alignment=PP_ALIGN.CENTER)


# ==================== SLIDE 2: A DIFFERENT KIND OF QUESTION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "A Different Kind of Question", font_size=28, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

add_card(slide, 0.8, 1.5, 5.6, 4.9)
tf = add_text(slide, 1.1, 1.7, 5.0, 0.5, "Every Strategy So Far Tried to Time Trades", font_size=18, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "A regime detector, a machine learning model, a", font_size=14, color=DARK)
add_para(tf, "cross-sectional ranking. All of them decide WHEN", font_size=14, color=DARK)
add_para(tf, "to buy and sell. All of them failed to beat the", font_size=14, color=DARK)
add_para(tf, "market consistently (Update 9).", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "So this update asks a structurally different", font_size=14, color=DARK)
add_para(tf, "question: what if we stop trying to predict", font_size=14, color=DARK)
add_para(tf, "direction at all?", font_size=14, color=DARK)

add_card(slide, 6.6, 1.5, 5.9, 4.9, fill_color=RGBColor(0xE4, 0xDE, 0xD2))
tf = add_text(slide, 6.9, 1.7, 5.3, 0.5, "Buy-and-Hold, With a Safety Net", font_size=18, color=DARK, bold=True)
add_para(tf, "")
add_para(tf, "The new strategy buys a stock once and holds it,", font_size=14, color=DARK)
add_para(tf, "exactly like a passive investor. It only sells if a", font_size=14, color=DARK)
add_para(tf, "stop-loss or trailing stop fires.", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "Let me explain what those are. A stop-loss sells", font_size=14, color=DARK)
add_para(tf, "automatically once the price falls a set percentage", font_size=14, color=DARK)
add_para(tf, "below where it was bought. A trailing stop instead", font_size=14, color=DARK)
add_para(tf, "sells once the price falls a set percentage below", font_size=14, color=DARK)
add_para(tf, "its highest point since being bought, so it locks in", font_size=14, color=DARK)
add_para(tf, "gains as a position rises.", font_size=14, color=DARK)
add_para(tf, "")
add_para(tf, "Once stopped out, it does not re-enter.", font_size=14, color=DARK, bold=True)

add_card(slide, 0.8, 6.55, 11.8, 0.7)
tf = add_text(slide, 1.1, 6.65, 11.3, 0.5, "This is the exact question Update 9 closed on: can risk management alone, with no signal at all, beat simply holding the stock?", font_size=13, color=MID)


# ==================== SLIDE 3: THE RESULTS ACROSS 34 STOCKS ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "The Results Across 34 Stocks", font_size=28, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

img = CHART_DIR / 'chart_u10_overview.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(2.0), Inches(1.5), width=Inches(9.3))

add_card(slide, 0.8, 6.15, 11.8, 1.15)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5, "The same two-configuration robust standard from every earlier study was applied here too.", font_size=14, color=DARK, bold=True)
add_para(tf, "Only 1 of 34 stocks passed it. 13 more beat plain buy-and-hold without passing it. 20 did neither.", font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 4: THE ONE ROBUST CASE STILL LOST ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "The One Robust Case Still Lost to Buy-and-Hold", font_size=26, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

img = CHART_DIR / 'chart_u10_jpm.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(2.7), Inches(1.5), width=Inches(8.0))

add_card(slide, 0.8, 6.15, 11.8, 1.15)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5, "JPMorgan (JPM) was the only stock where the overlay was profitable and consistent in both configurations.", font_size=14, color=DARK, bold=True)
add_para(tf, "Even so, it made less money than an investor who just bought JPM and did nothing at all.", font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 5: WHAT THE OTHER WINS ACTUALLY ARE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 0.5, 0.3, size=0.8, thickness=0.1)
add_text(slide, 0.8, 0.5, 12, 0.7, "What the Other 13 “Wins” Actually Are", font_size=27, color=DARK, bold=True)
add_divider(slide, 0.8, 1.15, 4)

img = CHART_DIR / 'chart_u10_lossmitigation.png'
if img.exists():
    slide.shapes.add_picture(str(img), Inches(1.7), Inches(1.5), width=Inches(9.9))

add_card(slide, 0.8, 6.15, 11.8, 1.15)
tf = add_text(slide, 1.1, 6.25, 11.3, 0.5, "These 5 stocks show the pattern most clearly: buy-and-hold itself was already losing money on all of them.", font_size=14, color=DARK, bold=True)
add_para(tf, "The overlay's stop-loss cut those losses short. It did not turn any of them into a profit.", font_size=14, color=MID, space_before=Pt(2))


# ==================== SLIDE 6: SUMMARY & NEXT ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_bracket_tl(slide, 1.5, 0.7, size=2.0, thickness=0.15)
add_bracket_br(slide, 11.8, 6.8, size=2.0, thickness=0.15)

add_text(slide, 2.5, 1.1, 8.5, 0.7, "SUMMARY", font_size=32, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)

tf = add_text(slide, 2.0, 2.1, 9.5, 0.5, "✅   Tested the one open question left from Update 9: risk management with no timing signal at all", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Same dual-configuration robust standard applied, across all 34 stocks used throughout this project", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Only 1 stock passed, and it still made less than doing nothing at all", font_size=17, color=DARK)
add_para(tf, "")
add_para(tf, "✅   Everywhere else, the overlay's only real effect was cutting losses shorter, never creating a gain", font_size=17, color=DARK)

add_card(slide, 2.0, 5.15, 9.5, 1.3)
tf = add_text(slide, 2.3, 5.3, 9, 0.5, "🎯  Next: Closing Out Cross-Sectional", font_size=18, color=ACCENT_GREEN, bold=True)
add_para(tf, "The research arc is reaching its final stage, but one thread is still open: cross-sectional reversal has so far only been tested on small, hand-picked baskets.", font_size=14, color=MID)
add_para(tf, "The next update tests it properly, on a full market index for real breadth, and on the opposite bet, momentum, before the report's findings are finalized.", font_size=14, color=MID, space_before=Pt(2))

add_text(slide, 2.5, 6.7, 8.5, 0.6, "Thank You", font_size=26, color=BRACKET, bold=True, alignment=PP_ALIGN.CENTER)


output_path = "/Users/jacksonetherchainstake/FYP/Documents/FYP_Update_10.pptx"
prs.save(output_path)
print(f"Saved to {output_path}")
